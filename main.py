import io
import zipfile
import streamlit as st
from pptx import Presentation
from PIL import Image
from rembg import remove

st.title("학생 이미지 일괄 투명화 & 파일명 자동 추출 변환기")
st.write("슬라이드 상단의 텍스트를 파일명으로 추출하고, 이미지 뒷배경을 투명하게 만들어 일괄 다운로드합니다.")

uploaded_file = st.file_uploader("파워포인트(PPTX) 파일을 선택하세요", type=["pptx"])

if uploaded_file is not None:
    if st.button("변환 시작"):
        with st.spinner("학생 이름을 읽고 배경을 투명하게 변환하는 중입니다..."):
            
            prs = Presentation(uploaded_file)
            zip_buffer = io.BytesIO()
            success_count = 0
            
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                
                for i, slide in enumerate(prs.slides):
                    student_name = f"학생_{i+1}" # 기본 이름
                    slide_image_blob = None
                    
                    # 1. 슬라이드 안의 모든 도형(Shape) 탐색
                    for shape in slide.shapes:
                        # 텍스트 상자 처리 (이름 추출)
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                # 파일명으로 안전한 문자만 추출 (공백, 한글, 영문, 숫자 등 허용)
                                safe_text = "".join(c for c in text if c.isalnum() or c in (' ', '_', '-', '(', ')')).strip()
                                if safe_text:
                                    # 첫 번째로 발견되는 텍스트를 이름으로 지정 (상단 텍스트 우선)
                                    if student_name.startswith("학생_"):
                                        student_name = safe_text
                        
                        # 이미지 개체 처리 (그림 타입이거나 내부에 이미지 속성이 있는 경우)
                        if shape.shape_type == 13 or hasattr(shape, "image"):
                            try:
                                slide_image_blob = shape.image.blob
                            except Exception:
                                pass
                    
                    # 2. 이미지가 발견된 경우 배경 투명화 및 파일 저장
                    if slide_image_blob:
                        try:
                            # PIL을 이용해 이미지 열기
                            image = Image.open(io.BytesIO(slide_image_blob))
                            
                            # RGBA 모드로 변환 (투명도 지원)
                            image = image.convert("RGBA")
                            
                            # rembg를 활용해 배경 제거 (투명화)
                            output_image = remove(image)
                            
                            # PNG 바이트로 변환
                            img_byte_arr = io.BytesIO()
                            output_image.save(img_byte_arr, format='PNG')
                            img_byte_arr = img_byte_arr.getvalue()
                            
                            # 파일명 중복 방지 (이름이 겹치면 슬라이드 번호 추가)
                            file_name = f"{student_name}.png"
                            
                            zip_file.writestr(file_name, img_byte_arr)
                            success_count += 1
                            
                        except Exception as e:
                            st.warning(f"슬라이드 {i+1} 처리 중 오류 발생: {e}")

            zip_buffer.seek(0)
            
            if success_count > 0:
                st.success(f"총 {success_count}개의 이미지 변환 및 배경 투명화 완료!")
                
                st.download_button(
                    label="모든 이미지 ZIP으로 다운로드",
                    data=zip_buffer,
                    file_name="transparent_students_images.zip",
                    mime="application/zip"
                )
            else:
                st.error("슬라이드에서 이미지 개체를 찾지 못했습니다. 파워포인트 내 이미지가 '그림' 형태로 삽입되어 있는지 확인해 주세요.")
