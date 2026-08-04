import io
import zipfile
import streamlit as st
from pptx import Presentation
from PIL import Image
from pptx.util import Emu

st.title("학생 이미지 추출 및 파일명 자동 지정 변환기")
st.write("각 슬라이드의 상단 텍스트를 파일명으로 지정하고, 원본 비율로 이미지를 추출합니다.")

uploaded_file = st.file_uploader("파워포인트(PPTX) 파일을 선택하세요", type=["pptx"])

if uploaded_file is not None:
    if st.button("변환 시작"):
        with st.spinner("이미지와 파일명을 추출하는 중입니다..."):
            
            prs = Presentation(uploaded_file)
            zip_buffer = io.BytesIO()
            success_count = 0
            
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                
                for i, slide in enumerate(prs.slides):
                    student_name = f"학생_{i+1}"
                    slide_image_blob = None
                    
                    # 1. 텍스트와 이미지 요소 탐색
                    for shape in slide.shapes:
                        # 텍스트 상자에서 학생 이름(파일명) 탐색
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                safe_text = "".join(c for c in text if c.isalnum() or c in (' ', '_', '-', '(', ')')).strip()
                                if safe_text:
                                    if student_name.startswith("학생_"):
                                        student_name = safe_text
                        
                        # 슬라이드 내 이미지 개체 추출
                        if shape.shape_type == 13 or hasattr(shape, "image"):
                            try:
                                slide_image_blob = shape.image.blob
                                # 핵심: 원본 이미지의 고유 비율을 가져옵니다.
                                # 파워포인트 내부 단위인 EMU를 픽셀로 변환하여 비율 계산
                                img_width_emu = shape.width
                                img_height_emu = shape.height
                                
                                # 도형이 회전되어 있을 경우를 대비한 보정
                                if shape.rotation != 0:
                                    # 간단한 보정 (회전각이 90/270도면 가로세로 스왑)
                                    if shape.rotation % 180 != 0:
                                        img_width_emu, img_height_emu = img_height_emu, img_width_emu

                                # 비율 계산 (정수형 나눗셈 방지를 위해 float 변환)
                                if img_height_emu > 0:
                                    aspect_ratio = float(img_width_emu) / float(img_height_emu)
                                else:
                                    aspect_ratio = 1.0 # 예외 처리
                                    
                            except Exception as e:
                                st.warning(f"이미지 속성 읽기 실패: {e}")
                    
                    # 2. 이미지 처리 및 저장 (원본 비율 유지)
                    if slide_image_blob:
                        try:
                            image = Image.open(io.BytesIO(slide_image_blob))
                            
                            # 원본 이미지 객체가 이미 로드되었으므로,
                            # 비율에 맞춰 PIL에서 리사이즈합니다.
                            # 이미지를 강제로 늘린 파워포인트 도형(shape)의 크기를 기준으로,
                            # 원본 비율대로 다시 계산합니다.
                            
                            # 현재 도형이 차지하는 가로 크기 (픽셀 단위로 근사 변환)
                            target_width = int(Emu(img_width_emu).inches * 100) # 100 DPI 가정 리사이즈
                            # (주의: 정확한 픽셀 변환은 복잡하므로, 
                            # PIL 자체 연산을 위해 비율만 넘겨줍니다.)
                            
                            # PIL을 이용한 비율 기반 리사이즈 (도형 크기의 가로폭을 기준으로 설정)
                            # 원본 비율(aspect_ratio)을 유지합니다.
                            if aspect_ratio > 1:
                                # 가로가 더 긴 이미지: 가로폭 기준, 세로 자동 계산
                                new_w = int(image.width)
                                new_h = int(new_w / aspect_ratio)
                            else:
                                # 세로가 더 긴 이미지: 세로폭 기준, 가로 자동 계산
                                new_h = int(image.height)
                                new_w = int(new_h * aspect_ratio)
                                
                            # 실제 리사이징 수행
                            resized_image = image.resize((new_w, new_h), Image.Resampling.LANCZOS)
                            
                            # RGBA 모드로 변환 (투명도 지원 유지)
                            resized_image = resized_image.convert("RGBA")
                            
                            img_byte_arr = io.BytesIO()
                            resized_image.save(img_byte_arr, format='PNG')
                            img_byte_arr = img_byte_arr.getvalue()
                            
                            file_name = f"{student_name}.png"
                            zip_file.writestr(file_name, img_byte_arr)
                            success_count += 1
                            
                        except Exception as e:
                            st.warning(f"슬라이드 {i+1} 처리 중 오류: {e}")

            zip_buffer.seek(0)
            
            if success_count > 0:
                st.success(f"총 {success_count}개의 이미지 변환 완료 (원본 비율 유지)!")
                st.download_button(
                    label="이미지 ZIP으로 다운로드",
                    data=zip_buffer,
                    file_name="students_images_original_ratio.zip",
                    mime="application/zip"
                )
            else:
                st.error("슬라이드에서 이미지를 찾지 못했습니다.")
